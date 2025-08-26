from flask import Flask, render_template, request, send_file, jsonify
import os, json, requests
from pptx import Presentation
from datetime import datetime
from pptx.util import Inches, Pt
from io import BytesIO

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'generated'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/parse', methods=['POST'])
def parse():
    text = request.form['text']
    guidance = request.form.get('guidance', '')
    api_key = request.form['api_key']
    api_url = request.form['api_url']
    model = request.form['model']

    # Prompt for LLM
    prompt = f"""
    You are an assistant that transforms raw text into PowerPoint presentation slides with minimal API usage.

    Task:
    1. Analyze the input text. If it is long, split it into contextually similar short sections.
    2. For each section, extract only the essential points.
    3. Generate one slide per section with:
       - A clear title
       - 3–5 concise bullet points
    4. The first slide should have only a suitable title derived from the full text topic.
    5. If a title is already present in the text, use that as the first slide.
    6. Output as JSON like:
       [
         {{"title": "Intro Slide Title", "bullets": []}},
         {{"title": "Section Title", "bullets": ["point 1", "point 2"]}}
       ]

    Rules:
    - Avoid copying long sentences from the input.
    - Do not include images or styling instructions.
    - Do not repeat information across slides.

    Input text:
    \"\"\"{guidance if guidance else ""}{text}\"\"\"
    """

    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    payload = {"model": model, "messages": [{"role": "user", "content": prompt}]}

    r = requests.post(api_url, headers=headers, json=payload)
    if not r.ok:
        return f"LLM request failed: {r.status_code} - {r.text}", 500

    response_content = r.json()['choices'][0]['message']['content']

    # --- JSON parse with error handling ---
    try:
        if "```" in response_content:  # Handle ```json wrapping
            response_content = response_content.split("```")[1]
            response_content = response_content.replace("json", "", 1).strip()

        parsed_slides = json.loads(response_content)
    except Exception as e:
        return f"Failed to parse JSON.<br><pre>{response_content}</pre><br>Error: {str(e)}", 500

    return render_template(
        "parsed_preview.html",
        parsed_slides=parsed_slides,
        text=text,
        guidance=guidance,
        api_key=api_key,
        api_url=api_url,
        model=model
    )


@app.route('/generate', methods=['POST'])
def generate():
    try:
        # 1. Read form data
        template_file = request.files['ppt_template']
        slides = json.loads(request.form['slides'])

        # 2. Load template
        prs = Presentation(template_file)

        # 3. Clear existing slides
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

        # 4. Find title layout
        title_layout = None
        for layout in prs.slide_layouts:
            if "title" in layout.name.lower():
                title_layout = layout
                break
        if not title_layout:
            title_layout = prs.slide_layouts[0]

        # 5. First slide: Title only
        first = prs.slides.add_slide(title_layout)
        if first.shapes.title:
            first.shapes.title.text = slides[0].get("title", "")
            p = first.shapes.title.text_frame.paragraphs[0]
            run = p.runs[0]
            run.font.bold = True
            run.font.size = Pt(36)

        # 6. Other slides: use same layout, add text boxes
        for s in slides[1:]:
            slide = prs.slides.add_slide(title_layout)

            # Subtitle box
            left, top, width, height = Inches(1), Inches(1), Inches(8), Inches(1)
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            tf = subtitle_box.text_frame
            tf.text = s.get("title", "")
            p = tf.paragraphs[0]
            run = p.runs[0]
            run.font.bold = True
            run.font.size = Pt(28)

            # Bullet box
            left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(5)
            bullet_box = slide.shapes.add_textbox(left, top, width, height)
            tf = bullet_box.text_frame
            for bullet in s.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.level = 0

        # 7. Return PPT
        output = BytesIO()
        prs.save(output)
        output.seek(0)

        return send_file(output, as_attachment=True, download_name="generated.pptx")

    except Exception as e:
        return f"Error in generate: {str(e)}", 500

if __name__ == '__main__':
    app.run(debug=True)
